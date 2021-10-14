from pptx import Presentation
prs = Presentation()
arquivo = input('nome do usuario e do arquivo:')
user = arquivo.split()[0]
nome = arquivo.split()[1]
#quantos slides
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
#configuração do slide 1
ti = input('titulo do texto:')
sub = input('subtitulo do texto:')
titulo = slide.shapes.title
titulo.text = ti
subtitulo = slide.placeholders[1]
subtitulo.text = sub
#configuração do slide 2
ti2 = input('topicos abordados:')
top1 = input('primeiro topico:')
top2 = input('segundo topico:')
top3 = input('terceiro topico:')
topicos = slide2.shapes.title
topicos.text = ti2
topico = slide2.placeholders[1]
topico.text = top1
topico2 = topico.text_frame.add_paragraph()
topico2.text = top2
topico3 = topico.text_frame.add_paragraph()
topico3.text = top3
prs.save('/Users/'+ user +'/Desktop/'+ nome +'.pptx')
#prs.save('pytest.pptx')
